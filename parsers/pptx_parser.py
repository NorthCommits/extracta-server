import logging
from typing import List, Dict
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from parsers.base_parser import BaseParser

logger = logging.getLogger("extracta.parser.pptx")

# EMU to points conversion (1 point = 12700 EMU)
EMU_TO_PT = 1 / 12700


class PPTXParser(BaseParser):

    def get_format(self) -> str:
        return "pptx"

    def extract_pages(self) -> List[Dict]:
        logger.info(f"[PPTXParser] Opening file: {self.file_path}")
        pages = []

        try:
            prs = Presentation(self.file_path)
        except Exception as e:
            logger.error(f"[PPTXParser] Failed to open file: {e}")
            raise

        slide_width = prs.slide_width * EMU_TO_PT
        slide_height = prs.slide_height * EMU_TO_PT

        logger.info(
            f"[PPTXParser] Total slides: {len(prs.slides)} "
            f"| Dimensions: {slide_width:.1f}x{slide_height:.1f} pt"
        )

        for slide_index, slide in enumerate(prs.slides):
            slide_number = slide_index + 1
            logger.debug(f"[PPTXParser] Processing slide {slide_number}")

            blocks = []

            # Process shapes in z-order (natural slide layer order)
            for shape in slide.shapes:
                shape_blocks = self._process_shape(
                    shape, slide_width, slide_height, slide_number
                )
                blocks.extend(shape_blocks)

            blocks = self.filter_empty_blocks(blocks)
            self.log_page_summary(slide_number, blocks)

            pages.append({
                "page_number": slide_number,
                "page_width": slide_width,
                "page_height": slide_height,
                "blocks": blocks,
            })

        logger.info(f"[PPTXParser] Extraction complete -- {len(pages)} slides processed")
        return pages

    def _process_shape(
        self,
        shape,
        slide_width: float,
        slide_height: float,
        slide_number: int
    ) -> List[dict]:
        """
        Process a single PPTX shape.
        Handles text boxes, placeholders, tables, and image shapes.
        """
        blocks = []

        try:
            bbox = self._get_shape_bbox(shape)
            if bbox is None:
                return blocks

            x0, y0, x1, y1 = bbox

            # Skip shapes outside slide bounds
            if x1 <= 0 or y1 <= 0 or x0 >= slide_width or y0 >= slide_height:
                logger.debug(f"[PPTXParser] Slide {slide_number} -- shape out of bounds, skipping")
                return blocks

            # Table shape
            if shape.has_table:
                table_block = self._process_table_shape(shape, bbox, slide_width, slide_height)
                if table_block:
                    blocks.append(table_block)
                return blocks

            # Text shape (text box, placeholder, auto shape with text)
            if shape.has_text_frame:
                text_blocks = self._process_text_shape(
                    shape, bbox, slide_width, slide_height
                )
                blocks.extend(text_blocks)
                return blocks

            # Image / picture shape
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blocks.append(self.make_block(
                    text="[IMAGE]",
                    bbox=bbox,
                    block_type="image",
                    page_width=slide_width,
                    page_height=slide_height,
                ))
                return blocks

        except Exception as e:
            logger.warning(f"[PPTXParser] Slide {slide_number} -- shape processing error: {e}")

        return blocks

    def _process_text_shape(
        self,
        shape,
        bbox: tuple,
        slide_width: float,
        slide_height: float
    ) -> List[dict]:
        """
        Extract text from a shape's text frame.
        Preserves paragraph structure and detects font properties.
        """
        results = []
        tf = shape.text_frame

        paragraphs_text = []
        font_sizes = []
        is_bold_flags = []

        for para in tf.paragraphs:
            para_text = ""
            for run in para.runs:
                para_text += run.text
                if run.font.size:
                    font_sizes.append(run.font.size * EMU_TO_PT)
                if run.font.bold:
                    is_bold_flags.append(True)

            # Also check paragraph-level text (no runs)
            if not para.runs and para.text:
                para_text = para.text

            if para_text.strip():
                paragraphs_text.append(para_text.strip())

        if not paragraphs_text:
            return results

        full_text = "\n".join(paragraphs_text)
        font_size = max(font_sizes) if font_sizes else None
        is_bold = any(is_bold_flags)

        # Check placeholder type for better classification
        placeholder_type = None
        if shape.is_placeholder:
            try:
                placeholder_type = shape.placeholder_format.type
            except Exception:
                pass

        block_type = self._classify_shape_block(
            shape, placeholder_type, font_size, is_bold, slide_height
        )

        results.append(self.make_block(
            text=full_text,
            bbox=bbox,
            block_type=block_type,
            font_size=font_size,
            is_bold=is_bold,
            page_width=slide_width,
            page_height=slide_height,
        ))

        return results

    def _process_table_shape(
        self,
        shape,
        bbox: tuple,
        slide_width: float,
        slide_height: float
    ) -> dict:
        """
        Extract table content as markdown-style text.
        """
        table = shape.table
        rows = []

        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                cells.append(cell_text)
            rows.append(" | ".join(cells))

        # Header separator
        if rows:
            header = rows[0]
            separator = " | ".join(["---"] * len(rows[0].split(" | ")))
            table_text = "\n".join([header, separator] + rows[1:])
        else:
            table_text = ""

        logger.debug(f"[PPTXParser] Table extracted with {len(rows)} rows")

        return self.make_block(
            text=table_text,
            bbox=bbox,
            block_type="table",
            page_width=slide_width,
            page_height=slide_height,
        )

    def _get_shape_bbox(self, shape) -> tuple:
        """
        Convert shape position and size from EMU to points.
        Returns (x0, y0, x1, y1) or None if position unavailable.
        """
        try:
            x0 = shape.left * EMU_TO_PT
            y0 = shape.top * EMU_TO_PT
            x1 = x0 + shape.width * EMU_TO_PT
            y1 = y0 + shape.height * EMU_TO_PT
            return self.normalize_bbox(x0, y0, x1, y1)
        except Exception as e:
            logger.debug(f"[PPTXParser] Could not get bbox for shape: {e}")
            return None

    def _classify_shape_block(
        self,
        shape,
        placeholder_type,
        font_size: float,
        is_bold: bool,
        slide_height: float
    ) -> str:
        """
        Classify shape block type using placeholder type, font size, and boldness.
        PPTX placeholder types:
            1 = CENTER_TITLE, 2 = BODY, 3 = CENTER_TITLE,
            13 = TITLE, 15 = SUBTITLE
        """
        from pptx.enum.text import PP_ALIGN

        # Use placeholder type if available
        if placeholder_type is not None:
            if placeholder_type in (1, 3, 13):
                return "title"
            elif placeholder_type == 15:
                return "heading"
            elif placeholder_type == 2:
                return "body"

        # Fallback to font heuristics
        if font_size is None:
            return "text"

        relative_size = font_size / slide_height if slide_height > 0 else 0

        if relative_size > 0.06 and is_bold:
            return "title"
        elif relative_size > 0.04:
            return "heading"
        elif relative_size < 0.015:
            return "footer"
        else:
            return "text"